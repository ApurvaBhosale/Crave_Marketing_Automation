# ======================================================
# Unified AI Content Generator (Blog + Video Script)
# ======================================================

import os
import fitz  # PyMuPDF #type: ignore
import docx #type: ignore
import requests #type: ignore
import warnings
import streamlit as st #type: ignore
from hdbcli import dbapi #type: ignore
from langchain_community.vectorstores.hanavector import HanaDB #type: ignore
from langchain_openai import AzureOpenAIEmbeddings #type: ignore
from langchain.chains import ConversationalRetrievalChain #type: ignore
#from langchain_community.document_loaders import TextLoade
from langchain.docstore.document import Document #type: ignore
from openai import AzureOpenAI #type: ignore
# for reading different file formats
from PyPDF2 import PdfReader #type: ignore
from docx import Document as DocxDocument #type: ignore
from pptx import Presentation #type: ignore

warnings.filterwarnings("ignore", category=DeprecationWarning)

# ======================================================
# Helper: Perplexity Search
# ======================================================
# Load Perplexity credentials from secrets.toml
PERPLEXITY_API_KEY = st.secrets["perplexity"]["api_key"]
PERPLEXITY_API_URL = st.secrets["perplexity"]["api_url"]

def perplexity_search(query, max_results=5):
    """Fetch results from Perplexity.ai"""
    payload = {"query": query}
    headers = {
        "Authorization": f"Bearer{PERPLEXITY_API_KEY}",
        "Content-Type": "application/json"
    }

    try:
        response = requests.post(PERPLEXITY_API_URL, json=payload, headers=headers, timeout=15)
        response.raise_for_status()
        data = response.json()
        results = []
        if "answer" in data:
            results.append(data["answer"])
        elif "data" in data and isinstance(data["data"], list):
            for item in data["data"][:max_results]:
                if "text" in item:
                    results.append(item["text"])
        return "\n".join(results)
    except Exception as e:
        return f"Perplexity API Error: {e}"

# ======================================================
# Helper: Read Documents
# ======================================================
def extract_text_from_file(file):
    """Extract text from PDF, DOCX, PPTX, or TXT"""
    text = ""
    name = file.name.lower()

    if name.endswith(".txt"):
        text = file.read().decode("utf-8")

    elif name.endswith(".pdf"):
        pdf = PdfReader(file)
        for page in pdf.pages:
            text += page.extract_text() or ""

    elif name.endswith(".docx"):
        doc = DocxDocument(file)
        text = "\n".join([p.text for p in doc.paragraphs])

    elif name.endswith(".pptx"):
        ppt = Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    return text.strip()

def extract_text_from_url(url):
    """Fetch and extract readable text from a given URL using Perplexity"""
    if not url or not url.strip():
        return ""

    try:
        # You can reuse your Perplexity API for this
        payload = {"query": f"Extract main article content from: {url}"}
        headers = {
            "Authorization": st.secrets["perplexity"]["api_key"],
            "Content-Type": "application/json"
        }
        response = requests.post(
            "https://api.perplexity.ai/search",
            json=payload,
            headers=headers,
            timeout=20
        )
        response.raise_for_status()
        data = response.json()
        if "answer" in data:
            return data["answer"]
        elif "data" in data and isinstance(data["data"], list):
            return "\n".join([d.get("text", "") for d in data["data"]])
    except Exception as e:
        return f"Error extracting URL content: {e}"

    return ""

# ======================================================
# Streamlit UI Setup
# ======================================================
st.set_page_config(page_title="AI Content Hub", layout="wide")
st.title("AI Content Hub")
# Subheading
st.markdown("##### AI-powered content creation for all your marketing needs")
# Add some vertical space before the next header
st.markdown("<br><br>", unsafe_allow_html=True)
#st.header("AI Content Generator")
# Row 1: Content Type
content_type = st.radio("**Select Content Type**", ["Blog", "Video Script"], horizontal=True)
# Row 2: Tone and Target Audience side by side
col1, col2 = st.columns(2)
with col1:
    tone = st.selectbox("**Select Tone**", ["Professional", "Friendly", "Authoritative", "Playful", "Inspirational"])
with col2:
    target_audience = st.selectbox(
        "**Target Audience**",
        ["Senior Management", "Middle Management", "Junior/Entry Level Staff"]
    )
# Row 3: Word Limit and Industry side by side
col3, col4 = st.columns(2)
with col3:
    word_limit = st.number_input("**Word Limit**", min_value=100, max_value=2000, value=1000)
with col4:
    industry = st.text_input("**Industry** (optional)")

# Row 5: Client name only for Video Script
#if content_type == "Video Script":
 #   client_name = st.text_input("Client Name (optional)")

# Row 6: Query input
query = st.text_input("**Enter your topic:**")
# Common field for both Blog and Video Script
additional_info = st.text_area(
    "**Add more information about the topic (optional):**",
    placeholder="Describe what you want to create. Be specific about your goals,etc...",
    height=120
)

uploaded_files = st.file_uploader(
    "**Upload reference document(s) (TXT, PDF, DOCX, PPTX)**",
    type=["txt", "pdf", "docx", "pptx"],
    accept_multiple_files=True
)


# Show uploaded files
if uploaded_files:
    st.write("**Uploaded files:**")
    for f in uploaded_files:
        st.write("📄", f.name)
        
##if content_type == "Blog":
  #  word_limit = st.number_input("Word Limit", min_value=100, max_value=2000, value=1000)
generate_button = st.button(f"Generate {content_type}")

# ======================================================
# Database + Azure Setup
# ======================================================
def init_services():
    connection = dbapi.connect(
        address=st.secrets["database"]["address"],
        port=st.secrets["database"]["port"],
        user=st.secrets["database"]["user"],
        password=st.secrets["database"]["password"],
        encrypt=True,
        autocommit=True,
        sslValidateCertificate=False,
    )

    client = AzureOpenAI(
        azure_endpoint=st.secrets["azure"]["openai_endpoint"],
        api_key=st.secrets["azure"]["api_key"],
        api_version=st.secrets["azure"]["api_version"],
    )

    embeddings = AzureOpenAIEmbeddings(
        azure_deployment=st.secrets["azure"]["embeddings_deployment"],
        openai_api_version=st.secrets["azure"]["embeddings_api_version"],
        api_key=st.secrets["azure"]["api_key"],
        azure_endpoint=st.secrets["azure"]["openai_endpoint"],
    )

    db = HanaDB(
        embedding=embeddings,
        connection=connection,
        table_name="MARKETING_APP_CONTENT_GENERATION"
    )
    return db, client

# ======================================================
# Unified RAG Logic
# ======================================================
def extract_text_from_url(url):
    if url.strip():
        return perplexity_search(url.strip())
    return ""

def retrieve_content(query, uploaded_files, db):
    """Retrieve content from HANA / Perplexity / Uploaded docs/URL"""
    
    # 1️⃣ Check uploaded files first
    uploaded_text = ""
    if uploaded_files:
        for f in uploaded_files:
            uploaded_text += extract_text_from_file(f) + "\n"
    if uploaded_text.strip():
        return uploaded_text 
    

    
    # Check HANA for company-specific content
    docs = db.similarity_search(query, k=20)
    hana_text = "\n".join([doc.page_content for doc in docs]) if docs else ""
    if hana_text.strip():
        return hana_text  # Use HANA content if it exists
    return ""

    


def generate_prompt_guidelines(tone, target_audience):
    """Return tone and audience-specific writing instructions"""
    # Tone-specific style guide
    tone_guidelines = {
        "Professional": "Use clear, concise, and confident language. Focus on credibility, precision, and business relevance.",
        "Friendly": "Use warm, conversational, and easy-to-understand language. Maintain professionalism but sound approachable.",
        "Authoritative": "Use confident, expert-driven language. Provide strong arguments and data-backed insights.",
        "Playful": "Use witty, light-hearted, and creative phrasing. Keep the tone fun yet informative, with clever transitions.",
        "Inspirational": "Use motivational and uplifting language. Focus on positive change, growth, and vision-driven storytelling."
    }
    tone_instruction = tone_guidelines.get(tone, "Use a balanced and clear writing tone suitable for professional readers.")

    # Audience-specific guidance
    audience_guidelines = {
        "Senior Management": (
            "Focus on strategic insights, ROI, and business impact. "
            "Use concise, high-level language. Avoid unnecessary technical details."
        ),
        "Middle Management": (
            "Provide actionable guidance, practical steps, and process-oriented insights. "
            "Balance strategic context with implementation advice."
        ),
        "Junior/Entry Level Staff": (
            "Explain clearly, use simple examples, and avoid jargon. "
            "Focus on learning, awareness, and foundational concepts."
        )
    }
    audience_instruction = audience_guidelines.get(target_audience, "")

    return tone_instruction, audience_instruction
# ======================================================
# Prompt & Generation Logic
# ======================================================
def generate_blog_prompt(tone, target_audience, industry, query, word_limit, final_content):
    tone_instruction, audience_instruction = generate_prompt_guidelines(tone, target_audience)
    return f"""
You are a skilled blog writer.

Write a {tone.lower()} blog for {target_audience}.

Tone: {tone}
Industry: {industry or "Not specified"}
Word Limit: {word_limit} words
Topic: {query}
If industry is specified, adapt examples, context, and recommendations to that industry.  
If industry is not specified, write for a general business/technology audience.  


**Strict Blog Structure and Formatting Guidelines:**
1.  **Title:**
    * Short, focused, and benefit-driven.
    * Must include important keywords relevant to the blog's topic.
2.  **Introduction (2 short paragraphs):**
    * Start with a compelling business pain point or a thought-provoking question.
    * Clearly explain the core problem addressed by the topic, using keywords naturally.
    * The second paragraph should seamlessly set the stage for what the reader will learn, without using phrases like "In this blog" or "This blog will cover." Instead, guide the reader into the content naturally.
    * Each paragraph should be 3-5 lines max.
3.  **Body:**
    * Use clear, descriptive subheadings for each distinct section.
    * After each subheading, provide a brief introductory paragraph (3-5 lines) that sets the context for that section.
    * Ensure a logical flow and strong, seamless connections between all sections and subheadings.
    * When presenting steps or points, ensure they are logically ordered and directly relevant to the surrounding content.
    * Explain technical concepts simply, but include sufficient depth where it adds significant value.
    * Integrate industry best practices, practical tips, or relevant SAP standards.
    * **General Body Formatting:**
        * Each paragraph must be short (3-5 lines).
        * Each sentence must be 20-25 words maximum.
        * Use bullet points or numbered lists where they enhance readability and clarity.
4.  **Conclusion**
    * **Heading:** Create a concise, appropriate, and engaging heading for this final section. Avoid using "Wrap-up," "Summary," or "Call to Action" directly in the heading itself. Think of a heading that summarizes the benefit or next step.
    * Provide a brief, concise summary of the key takeaways presented in the blog.
 
**Prohibited Content & Phrases:**
* Do NOT use generic introductory phrases such as "In today's fast-paced," "In this blog," "This blog will cover," "We will explore," or similar.
* Do NOT generate generic, fabricated, or made-up blog content. All content must stem from real insights and the provided context.
* Do NOT include any testing instructions or conversational filler in the final blog output.



Use the following reference content:
{final_content}
"""

def generate_video_prompt(tone, target_audience, industry, final_content):
    tone_instruction, audience_instruction = generate_prompt_guidelines(tone, target_audience)
    return f"""
You are an expert video scriptwriter.

Write a {tone.lower()} video script for {target_audience} audience.
Focus on the {industry or "Not specified"} industry.

Structure:
1. Problem Introduction
2. Product/Brand Introduction
3. Key Features Highlights
4. Benefit Explanantion
5. Real Life Example or Case Study
6. Call-to-Action
7. Closing Scene

Each scene should include:

- Visual: On-screen visuals
- Narration: Voiceover

Use the following reference content:
{final_content}
"""

# ======================================================
# Generate Content
# ======================================================
if generate_button and query:
    with st.spinner(f"Generating {content_type}..."):
        db, client = init_services()
        
        # Combine query + additional info
        full_query = query
        if additional_info.strip():
            full_query += f"\n\nAdditional Information:\n{additional_info.strip()}"

        final_content = retrieve_content(full_query, uploaded_files, db)
        

        if content_type == "Blog":
            prompt = generate_blog_prompt(tone, target_audience, industry, full_query, word_limit, final_content)
        else:
            prompt = generate_video_prompt(tone, target_audience, industry, final_content)

        messages = [{"role": "system", "content": prompt}]
        response = client.chat.completions.create(
            messages=messages,
            model="Codetest",
            max_tokens=1600,
            temperature=0.7
        )

        output = response.choices[0].message.content

    st.subheader(f"Generated {content_type} ✨")
    st.markdown(output)





